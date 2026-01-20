#!/usr/bin/env python3
"""
Utilitaires pour la manipulation de documents.
PDF, DOCX, PPTX, XLSX, images, LaTeX, etc.

Auteur: [Votre nom]
Memoire: Mesures DNS dans l'espace et le temps
"""

from pathlib import Path
from typing import Union, List, Optional
import subprocess
import tempfile

# === PDF ===
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
import pdfplumber
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

# === Documents Office ===
from docx import Document
from pptx import Presentation
from openpyxl import Workbook, load_workbook

# === Images ===
from PIL import Image

# === Markdown ===
import markdown

# === Logging ===
from loguru import logger


# ============================================================================
# PDF UTILITIES
# ============================================================================

def pdf_to_text(pdf_path: Union[str, Path]) -> str:
    """
    Extrait le texte d'un fichier PDF.

    Args:
        pdf_path: Chemin vers le fichier PDF

    Returns:
        Texte extrait
    """
    text_parts = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)

    return '\n\n'.join(text_parts)


def pdf_get_info(pdf_path: Union[str, Path]) -> dict:
    """
    Recupere les metadonnees d'un PDF.

    Args:
        pdf_path: Chemin vers le fichier PDF

    Returns:
        Dictionnaire avec les informations
    """
    reader = PdfReader(pdf_path)
    info = reader.metadata

    return {
        'pages': len(reader.pages),
        'title': info.get('/Title', None) if info else None,
        'author': info.get('/Author', None) if info else None,
        'subject': info.get('/Subject', None) if info else None,
        'creator': info.get('/Creator', None) if info else None
    }


def merge_pdfs(pdf_paths: List[Union[str, Path]], output_path: Union[str, Path]):
    """
    Fusionne plusieurs fichiers PDF.

    Args:
        pdf_paths: Liste des chemins PDF
        output_path: Chemin du fichier de sortie
    """
    merger = PdfMerger()
    for pdf in pdf_paths:
        merger.append(str(pdf))
    merger.write(str(output_path))
    merger.close()
    logger.success(f"PDFs fusionnes: {output_path}")


def split_pdf(pdf_path: Union[str, Path], output_dir: Union[str, Path]):
    """
    Separe un PDF en pages individuelles.

    Args:
        pdf_path: Chemin vers le PDF
        output_dir: Repertoire de sortie
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(exist_ok=True)

    reader = PdfReader(pdf_path)
    for i, page in enumerate(reader.pages):
        writer = PdfWriter()
        writer.add_page(page)

        output_file = output_dir / f"page_{i+1:03d}.pdf"
        with open(output_file, 'wb') as f:
            writer.write(f)

    logger.success(f"PDF separe en {len(reader.pages)} fichiers dans {output_dir}")


def pdf_extract_tables(pdf_path: Union[str, Path]) -> list:
    """
    Extrait les tableaux d'un PDF.

    Args:
        pdf_path: Chemin vers le PDF

    Returns:
        Liste de tableaux (listes de listes)
    """
    tables = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_tables = page.extract_tables()
            tables.extend(page_tables)

    return tables


# ============================================================================
# WORD (DOCX) UTILITIES
# ============================================================================

def docx_to_text(docx_path: Union[str, Path]) -> str:
    """
    Extrait le texte d'un fichier Word.

    Args:
        docx_path: Chemin vers le fichier DOCX

    Returns:
        Texte extrait
    """
    doc = Document(docx_path)
    text_parts = [paragraph.text for paragraph in doc.paragraphs]
    return '\n'.join(text_parts)


def create_docx(output_path: Union[str, Path], title: str, content: str):
    """
    Cree un document Word simple.

    Args:
        output_path: Chemin de sortie
        title: Titre du document
        content: Contenu textuel
    """
    doc = Document()
    doc.add_heading(title, 0)

    for paragraph in content.split('\n\n'):
        doc.add_paragraph(paragraph)

    doc.save(output_path)
    logger.success(f"Document Word cree: {output_path}")


def docx_to_pdf(docx_path: Union[str, Path], output_path: Optional[Union[str, Path]] = None):
    """
    Convertit un DOCX en PDF (via LibreOffice).

    Args:
        docx_path: Chemin vers le DOCX
        output_path: Chemin de sortie (optionnel)
    """
    docx_path = Path(docx_path)
    if output_path is None:
        output_path = docx_path.with_suffix('.pdf')

    output_dir = Path(output_path).parent

    cmd = [
        'libreoffice',
        '--headless',
        '--convert-to', 'pdf',
        '--outdir', str(output_dir),
        str(docx_path)
    ]

    subprocess.run(cmd, check=True)
    logger.success(f"DOCX converti en PDF: {output_path}")


# ============================================================================
# POWERPOINT (PPTX) UTILITIES
# ============================================================================

def pptx_to_text(pptx_path: Union[str, Path]) -> str:
    """
    Extrait le texte d'une presentation PowerPoint.

    Args:
        pptx_path: Chemin vers le fichier PPTX

    Returns:
        Texte extrait
    """
    prs = Presentation(pptx_path)
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"=== Slide {slide_num} ==="]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text_parts.append('\n'.join(slide_text))

    return '\n\n'.join(text_parts)


def create_pptx(output_path: Union[str, Path], slides_content: List[dict]):
    """
    Cree une presentation PowerPoint.

    Args:
        output_path: Chemin de sortie
        slides_content: Liste de dicts avec 'title' et 'content'
    """
    prs = Presentation()

    for slide_data in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = slide_data.get('title', '')
        slide.placeholders[1].text = slide_data.get('content', '')

    prs.save(output_path)
    logger.success(f"Presentation PowerPoint creee: {output_path}")


# ============================================================================
# EXCEL (XLSX) UTILITIES
# ============================================================================

def xlsx_to_dataframes(xlsx_path: Union[str, Path]) -> dict:
    """
    Lit un fichier Excel et retourne un dict de DataFrames.

    Args:
        xlsx_path: Chemin vers le fichier XLSX

    Returns:
        Dict avec nom de feuille -> DataFrame
    """
    import pandas as pd
    return pd.read_excel(xlsx_path, sheet_name=None)


def dataframe_to_xlsx(df, output_path: Union[str, Path], sheet_name: str = 'Sheet1'):
    """
    Sauvegarde un DataFrame en fichier Excel.

    Args:
        df: pandas DataFrame
        output_path: Chemin de sortie
        sheet_name: Nom de la feuille
    """
    df.to_excel(output_path, sheet_name=sheet_name, index=False)
    logger.success(f"Excel cree: {output_path}")


# ============================================================================
# IMAGE UTILITIES
# ============================================================================

def image_info(image_path: Union[str, Path]) -> dict:
    """
    Recupere les informations d'une image.

    Args:
        image_path: Chemin vers l'image

    Returns:
        Dict avec les informations
    """
    img = Image.open(image_path)
    return {
        'format': img.format,
        'mode': img.mode,
        'size': img.size,
        'width': img.width,
        'height': img.height
    }


def convert_image(
    input_path: Union[str, Path],
    output_path: Union[str, Path],
    size: Optional[tuple] = None
):
    """
    Convertit une image vers un autre format et/ou redimensionne.

    Args:
        input_path: Chemin source
        output_path: Chemin destination
        size: Nouvelle taille (width, height) optionnelle
    """
    img = Image.open(input_path)

    if size:
        img = img.resize(size, Image.Resampling.LANCZOS)

    # Convertir en RGB si necessaire pour JPEG
    output_path = Path(output_path)
    if output_path.suffix.lower() in ['.jpg', '.jpeg'] and img.mode in ['RGBA', 'P']:
        img = img.convert('RGB')

    img.save(output_path)
    logger.success(f"Image convertie: {output_path}")


def images_to_pdf(image_paths: List[Union[str, Path]], output_path: Union[str, Path]):
    """
    Combine plusieurs images en un seul PDF.

    Args:
        image_paths: Liste des chemins d'images
        output_path: Chemin du PDF de sortie
    """
    images = []
    for img_path in image_paths:
        img = Image.open(img_path)
        if img.mode == 'RGBA':
            img = img.convert('RGB')
        images.append(img)

    if images:
        images[0].save(
            output_path,
            save_all=True,
            append_images=images[1:],
            resolution=100.0
        )
        logger.success(f"Images combinees en PDF: {output_path}")


# ============================================================================
# LATEX UTILITIES
# ============================================================================

def compile_latex(
    tex_path: Union[str, Path],
    output_dir: Optional[Union[str, Path]] = None,
    engine: str = 'pdflatex'
):
    """
    Compile un fichier LaTeX en PDF.

    Args:
        tex_path: Chemin vers le fichier .tex
        output_dir: Repertoire de sortie (optionnel)
        engine: Moteur LaTeX (pdflatex, xelatex, lualatex)
    """
    tex_path = Path(tex_path)
    if output_dir is None:
        output_dir = tex_path.parent

    cmd = [
        'latexmk',
        f'-{engine}',
        '-interaction=nonstopmode',
        f'-output-directory={output_dir}',
        str(tex_path)
    ]

    subprocess.run(cmd, check=True, cwd=tex_path.parent)
    logger.success(f"LaTeX compile: {tex_path.stem}.pdf")


def markdown_to_latex(md_content: str) -> str:
    """
    Convertit du Markdown en LaTeX (basique).

    Args:
        md_content: Contenu Markdown

    Returns:
        Contenu LaTeX
    """
    # Conversion basique via pandoc
    result = subprocess.run(
        ['pandoc', '-f', 'markdown', '-t', 'latex'],
        input=md_content,
        capture_output=True,
        text=True
    )
    return result.stdout


def markdown_to_pdf(md_path: Union[str, Path], output_path: Union[str, Path]):
    """
    Convertit un fichier Markdown en PDF via pandoc.

    Args:
        md_path: Chemin vers le fichier Markdown
        output_path: Chemin du PDF de sortie
    """
    cmd = [
        'pandoc',
        str(md_path),
        '-o', str(output_path),
        '--pdf-engine=xelatex'
    ]

    subprocess.run(cmd, check=True)
    logger.success(f"Markdown converti en PDF: {output_path}")


# ============================================================================
# MAIN - Demonstration
# ============================================================================

if __name__ == '__main__':
    logger.info("=== Utilitaires de documents ===")

    # Exemple: creer un document Word
    create_docx(
        '/workspace/output/test_document.docx',
        'Test Document',
        'Ceci est un paragraphe de test.\n\nEt voici un autre paragraphe.'
    )

    # Exemple: creer une presentation
    slides = [
        {'title': 'Introduction', 'content': 'Bienvenue dans cette presentation'},
        {'title': 'Contenu', 'content': 'Voici le contenu principal'},
        {'title': 'Conclusion', 'content': 'Merci de votre attention'}
    ]
    create_pptx('/workspace/output/test_presentation.pptx', slides)

    logger.success("Demonstrations terminees!")
